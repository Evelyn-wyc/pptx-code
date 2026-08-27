#!/usr/bin/env python3
"""Run deterministic font, notes, capacity, and overlap checks on a PPTX."""

from __future__ import annotations

import argparse
import json
import math
import sys
import unicodedata
from collections.abc import Iterable
from dataclasses import asdict, dataclass
from pathlib import Path

try:
    from pptx import Presentation
except ImportError as exc:  # pragma: no cover - exercised only without dependency
    raise SystemExit("python-pptx is required: pip install python-pptx") from exc

EMU_PER_INCH = 914_400
FOOTER_TOKENS = ("FOOTER", "SOURCE", "CITATION", "PAGE_NUMBER", "SLIDE_NUMBER")


@dataclass(frozen=True)
class Finding:
    slide: int
    kind: str
    item: str
    detail: str


@dataclass(frozen=True)
class TextItem:
    name: str
    text: str
    left: int
    top: int
    width: int
    height: int
    font_sizes: tuple[float, ...]

    @property
    def right(self) -> int:
        return self.left + self.width

    @property
    def bottom(self) -> int:
        return self.top + self.height


def parser() -> argparse.ArgumentParser:
    cli = argparse.ArgumentParser(description=__doc__)
    cli.add_argument("pptx", type=Path, help="PPTX file to inspect")
    cli.add_argument("--min-font", type=float, default=22.0)
    cli.add_argument("--require-notes", action="store_true")
    cli.add_argument("--skip-overlaps", action="store_true")
    cli.add_argument("--skip-capacity", action="store_true")
    cli.add_argument("--json", action="store_true", dest="as_json")
    return cli


def text_units(text: str) -> float:
    units = 0.0
    for char in text:
        if char in "\r\n":
            continue
        units += 1.0 if unicodedata.east_asian_width(char) in {"W", "F"} else 0.55
    return units


def explicit_font_sizes(frame) -> tuple[float, ...]:
    sizes: list[float] = []
    for paragraph in frame.paragraphs:
        for run in paragraph.runs:
            if run.text.strip() and run.font.size is not None:
                sizes.append(round(run.font.size.pt, 2))
    return tuple(sizes)


def iter_text_items(slide) -> Iterable[TextItem]:
    for shape_index, shape in enumerate(slide.shapes, start=1):
        shape_name = shape.name or f"shape-{shape_index}"
        if getattr(shape, "has_text_frame", False):
            text = shape.text_frame.text.strip()
            if text:
                yield TextItem(
                    shape_name,
                    text,
                    shape.left,
                    shape.top,
                    shape.width,
                    shape.height,
                    explicit_font_sizes(shape.text_frame),
                )
        if not getattr(shape, "has_table", False):
            continue
        table = shape.table
        y = shape.top
        for row_index, row in enumerate(table.rows, start=1):
            x = shape.left
            for column_index, column in enumerate(table.columns, start=1):
                cell = table.cell(row_index - 1, column_index - 1)
                text = cell.text_frame.text.strip()
                if text:
                    yield TextItem(
                        f"{shape_name}[r{row_index}c{column_index}]",
                        text,
                        x,
                        y,
                        column.width,
                        row.height,
                        explicit_font_sizes(cell.text_frame),
                    )
                x += column.width
            y += row.height


def is_footer(item: TextItem, slide_height: int) -> bool:
    upper_name = item.name.upper()
    named_footer = any(token in upper_name for token in FOOTER_TOKENS)
    bottom_band = (
        item.top >= slide_height * 0.925 and item.height <= slide_height * 0.075
    )
    return named_footer or bottom_band


def capacity_finding(slide_number: int, item: TextItem) -> Finding | None:
    if not item.text or not item.font_sizes:
        return None
    font_size = max(item.font_sizes)
    width_pt = item.width / EMU_PER_INCH * 72
    height_pt = item.height / EMU_PER_INCH * 72
    usable_width = max(width_pt - font_size * 0.35, font_size)
    usable_height = max(height_pt - font_size * 0.25, font_size)
    line_capacity = max(usable_width / (font_size * 0.92), 1.0)
    estimated_lines = 0
    for raw_line in item.text.splitlines() or [item.text]:
        estimated_lines += max(1, math.ceil(text_units(raw_line) / line_capacity))
    height_capacity = max(usable_height / (font_size * 1.12), 1.0)
    if estimated_lines <= height_capacity + 0.75:
        return None
    return Finding(
        slide_number,
        "capacity",
        item.name,
        f"estimated {estimated_lines} lines in space for {height_capacity:.1f}",
    )


def overlap_ratio(first: TextItem, second: TextItem) -> tuple[float, float, float]:
    width = max(0, min(first.right, second.right) - max(first.left, second.left))
    height = max(0, min(first.bottom, second.bottom) - max(first.top, second.top))
    if width == 0 or height == 0:
        return 0.0, 0.0, 0.0
    area = width * height
    smaller = min(first.width * first.height, second.width * second.height)
    return area / smaller, width / EMU_PER_INCH, height / EMU_PER_INCH


def validate(
    path: Path, min_font: float, require_notes: bool, overlaps: bool, capacity: bool
) -> list[Finding]:
    deck = Presentation(path)
    findings: list[Finding] = []
    for slide_number, slide in enumerate(deck.slides, start=1):
        items = list(iter_text_items(slide))
        for item in items:
            if not is_footer(item, deck.slide_height):
                for size in item.font_sizes:
                    if size + 0.01 < min_font:
                        findings.append(
                            Finding(
                                slide_number,
                                "font",
                                item.name,
                                f"{size:g} pt is below {min_font:g} pt",
                            )
                        )
                        break
            if capacity:
                finding = capacity_finding(slide_number, item)
                if finding is not None:
                    findings.append(finding)
        if overlaps:
            shape_items = [item for item in items if "[r" not in item.name]
            for index, first in enumerate(shape_items):
                for second in shape_items[index + 1 :]:
                    ratio, width_inches, height_inches = overlap_ratio(first, second)
                    if ratio >= 0.18 and width_inches >= 0.08 and height_inches >= 0.08:
                        findings.append(
                            Finding(
                                slide_number,
                                "overlap",
                                f"{first.name} <> {second.name}",
                                f"intersection is {ratio:.0%} of the smaller box",
                            )
                        )
        if require_notes:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if len(notes) < 20:
                findings.append(
                    Finding(
                        slide_number,
                        "notes",
                        "speaker notes",
                        "missing or shorter than 20 characters",
                    )
                )
    return findings


def main() -> int:
    args = parser().parse_args()
    path = args.pptx.resolve()
    if not path.is_file():
        raise SystemExit(f"PPTX not found: {path}")
    findings = validate(
        path,
        min_font=args.min_font,
        require_notes=args.require_notes,
        overlaps=not args.skip_overlaps,
        capacity=not args.skip_capacity,
    )
    if args.as_json:
        print(
            json.dumps(
                [asdict(item) for item in findings], ensure_ascii=False, indent=2
            )
        )
    elif findings:
        for finding in findings:
            print(
                f"slide {finding.slide}: {finding.kind}: {finding.item}: {finding.detail}"
            )
    else:
        print(f"OK: {path}")
    return 1 if findings else 0


if __name__ == "__main__":
    sys.exit(main())
